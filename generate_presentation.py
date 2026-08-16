# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

# ── Colors ──────────────────────────────
DARK  = RGBColor(0x06,0x0A,0x1A)   # near-black navy
MID   = RGBColor(0x0F,0x24,0x5E)   # deep blue
ACCT  = RGBColor(0xC9,0xA0,0x2C)   # gold
CYAN  = RGBColor(0x00,0xD4,0xFF)   # electric cyan
PRPL  = RGBColor(0x7B,0x2F,0xFF)   # purple accent
WHITE = RGBColor(0xFF,0xFF,0xFF)
LGRAY = RGBColor(0xE8,0xEC,0xF8)

def prs_init():
    p = Presentation()
    p.slide_width  = Inches(13.33)
    p.slide_height = Inches(7.5)
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, l, t, w, h, color, alpha=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def circle(slide, l, t, sz, color):
    s = slide.shapes.add_shape(9, Inches(l), Inches(t), Inches(sz), Inches(sz))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb

def mlabel(slide, lines, l, t, w, h, size=15, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(3)
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Calibri"

def transition(slide, t="fade"):
    xml = f'<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:{t}/></p:transition>'
    slide._element.insert(-1, etree.fromstring(xml))

def section_slide(prs, num, title, color_a=CYAN, color_b=PRPL):
    s = blank(prs)
    rect(s,  0,  0, 13.33, 7.5, DARK)
    # Big background number
    rect(s, 6.5, 0.5,  6.5, 6.5, MID)
    circle(s, 7.0, 0.0, 7.0, RGBColor(0x10,0x20,0x50))
    # Cyan glow bar left
    rect(s, 0, 0, 0.35, 7.5, color_a)
    rect(s, 0.35, 0, 0.08, 7.5, RGBColor(0x00,0x80,0xAA))
    # Gold bottom bar
    rect(s, 0, 7.1, 13.33, 0.4, ACCT)
    txt(s, num,   0.7, 0.6, 5, 3.5,  size=140, bold=True, color=RGBColor(0x15,0x30,0x70))
    txt(s, title, 0.7, 4.8, 11, 1.5, size=38, bold=True,  color=WHITE)
    rect(s, 0.7, 4.6, 4.0, 0.05, color_a)
    transition(s, "fade")
    return s

def header_slide(prs, title):
    s = blank(prs)
    rect(s, 0, 0, 13.33, 7.5, LGRAY)
    rect(s, 0, 0, 13.33, 1.35, DARK)
    # Decorative dots top-right
    for i in range(5):
        circle(s, 11.5+i*0.35, 0.1, 0.22, MID)
    txt(s, title, 0.5, 0.22, 12, 0.9, size=27, bold=True, color=WHITE)
    rect(s, 0.5, 1.45, 12.2, 0.05, ACCT)
    return s

# ─────────────────────────────────────────
prs = prs_init()

# ══════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════
s = blank(prs)
rect(s, 0, 0, 13.33, 7.5, DARK)
rect(s, 0, 0, 0.4,  7.5, CYAN)
rect(s, 0.4, 0, 0.1, 7.5, RGBColor(0,100,140))
circle(s, 9.5, -1.2, 6.0, MID)
circle(s, 10.2,-0.4, 4.0, RGBColor(0x10,0x20,0x60))
circle(s,-1.5, 5.0,  4.0, RGBColor(0x10,0x20,0x60))
rect(s, 0, 7.1, 13.33, 0.4, ACCT)
rect(s, 0.7, 1.0, 7.5, 0.05, CYAN)
txt(s, "BTS DEVELOPPEMENT SUR INTERNET", 0.7, 0.5, 10, 0.5, size=11, color=CYAN, bold=True)
txt(s, "LexPro", 0.7, 1.1, 10, 2.2, size=95, bold=True, color=WHITE)
txt(s, "Conception et developpement d'une application web\nde gestion pour cabinet d'avocats",
    0.7, 3.3, 10, 1.3, size=19, color=ACCT)
rect(s, 0.7, 4.75, 6, 0.05, RGBColor(0x40,0x50,0x80))
txt(s, "Elabore par :  Med Khalil Kribi",   0.7, 4.9, 8, 0.5, size=16, color=WHITE, bold=True)
txt(s, "Encadre par : M. Walid Ben Abda",   0.7, 5.45, 8, 0.5, size=15, color=LGRAY)
txt(s, "2024 / 2025", 0.7, 6.05, 5, 0.4, size=13, color=ACCT, italic=True)
transition(s, "fade")

# ══════════════════════════════════════════
# SLIDE 2 — Plan (6 cards)
# ══════════════════════════════════════════
s = blank(prs)
rect(s, 0, 0, 13.33, 7.5, DARK)
rect(s, 0, 0, 13.33, 0.08, CYAN)
rect(s, 0, 7.42, 13.33, 0.08, ACCT)
txt(s, "PLAN", 0, 0.15, 13.33, 0.9, size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rect(s, 4.5, 1.05, 4.33, 0.05, ACCT)

plan = [("01","Problematique"),("02","Specification"),("03","Conception"),
        ("04","Technologies"),("05","Realisation"),("06","Conclusion")]
pos  = [(0.3,1.6),(4.5,1.6),(8.7,1.6),(0.3,4.2),(4.5,4.2),(8.7,4.2)]

for i,((num,lbl),(cx,cy)) in enumerate(zip(plan,pos)):
    c = slide = s
    bg = MID if i%2==0 else RGBColor(0x12,0x1A,0x40)
    rect(c, cx, cy, 3.9, 2.55, bg)
    rect(c, cx, cy, 0.12, 2.55, CYAN if i<3 else ACCT)
    circle(c, cx+2.8, cy-0.35, 1.3, RGBColor(0x0A,0x18,0x50))
    txt(c, num,  cx+0.22, cy+0.18, 2.0, 0.8, size=32, bold=True, color=CYAN if i<3 else ACCT)
    txt(c, lbl,  cx+0.22, cy+1.1,  3.5, 1.1, size=15, color=WHITE, wrap=True)

transition(s, "push")

# ══════════════════════════════════════════
# SLIDE 3 — Section Problematique
# ══════════════════════════════════════════
section_slide(prs, "01", "Problematique", CYAN, PRPL)

# ══════════════════════════════════════════
# SLIDE 4 — Problematique detail
# ══════════════════════════════════════════
s = header_slide(prs, "Problematique")
# Left box
rect(s, 0.4, 1.6, 5.8, 5.5, MID)
rect(s, 0.4, 1.6, 0.12, 5.5, CYAN)
txt(s, "Contexte", 0.65, 1.75, 5, 0.55, size=17, bold=True, color=CYAN)
rect(s, 0.65, 2.35, 5.3, 0.04, CYAN)
txt(s, "Gerer un cabinet d'avocats implique de manipuler des centaines de dossiers, rendez-vous et documents disperses, causant des pertes de temps importantes et des risques d'erreurs.",
    0.65, 2.5, 5.3, 2.5, size=14, color=WHITE, wrap=True)
# Right box
rect(s, 7.0, 1.6, 6.0, 5.5, RGBColor(0xED,0xF0,0xFC))
rect(s, 7.0, 1.6, 0.12, 5.5, ACCT)
txt(s, "Problematique centrale", 7.25, 1.75, 5.5, 0.55, size=16, bold=True, color=MID)
rect(s, 7.25, 2.35, 5.5, 0.04, ACCT)
txt(s, "Comment centraliser sur une seule plateforme web securisee :\n\n- La gestion des dossiers juridiques\n- Le suivi du temps facturable\n- La facturation automatisee\n- L'agenda et les rendez-vous\n\n...pour les avocats et leurs assistants ?",
    7.25, 2.5, 5.5, 4.3, size=14, color=MID, wrap=True)
# Divider
rect(s, 6.45, 1.6, 0.1, 5.5, ACCT)
transition(s, "push")

# ══════════════════════════════════════════
# SLIDE 5 — Section Specification
# ══════════════════════════════════════════
section_slide(prs, "02", "Specification des Besoins", ACCT, CYAN)

# ══════════════════════════════════════════
# SLIDE 6 — Goals (5 cards)
# ══════════════════════════════════════════
s = header_slide(prs, "But du Projet")
goals = [
    ("Centraliser","Regrouper clients, dossiers et documents en un seul espace."),
    ("Optimiser","Time-Tracking et suivi du temps facturable integre."),
    ("Simplifier","Facturation et calcul des honoraires automatises."),
    ("Organiser","Agenda interactif et gestion des rendez-vous."),
    ("Securiser","Controle d'acces RBAC par profil (Avocat / Assistant)."),
]
colors = [CYAN, ACCT, PRPL, RGBColor(0x00,0xC8,0x80), RGBColor(0xFF,0x6B,0x35)]
for i, ((title, desc), color) in enumerate(zip(goals, colors)):
    cx = 0.28 + i*2.57
    rect(s, cx, 1.6, 2.35, 5.55, MID)
    rect(s, cx, 1.6, 2.35, 0.18, color)
    txt(s, str(i+1).zfill(2), cx+0.1, 1.85, 0.8, 0.8, size=28, bold=True, color=color)
    txt(s, title, cx+0.1, 2.7, 2.1, 0.7, size=14, bold=True, color=color, wrap=True, align=PP_ALIGN.CENTER)
    txt(s, desc,  cx+0.1, 3.55, 2.1, 3.3, size=12, color=WHITE, wrap=True, align=PP_ALIGN.CENTER)
transition(s, "push")

# ══════════════════════════════════════════
# SLIDE 7 — Section Conception
# ══════════════════════════════════════════
section_slide(prs, "03", "Conception", PRPL, CYAN)

# ══════════════════════════════════════════
# SLIDE 8 — Cas d'utilisation
# ══════════════════════════════════════════
s = header_slide(prs, "Conception  |  Diagramme de Cas d'utilisation")
rect(s, 0.5, 1.6, 12.3, 5.55, RGBColor(0xD8,0xDF,0xF5))
rect(s, 0.5, 1.6, 12.3, 0.1,  PRPL)
txt(s, "[ Inserer ici la capture du Diagramme de Cas d'utilisation ]",
    0.6, 3.7, 12, 1.0, size=17, color=RGBColor(0x55,0x60,0x80), align=PP_ALIGN.CENTER, bold=True)
transition(s, "push")

# ══════════════════════════════════════════
# SLIDE 9 — Diagramme de Classe
# ══════════════════════════════════════════
s = header_slide(prs, "Conception  |  Diagramme de Classe")
rect(s, 0.5, 1.6, 12.3, 5.55, RGBColor(0xD8,0xDF,0xF5))
rect(s, 0.5, 1.6, 12.3, 0.1,  CYAN)
txt(s, "[ Inserer ici la capture du Diagramme de Classe ]",
    0.6, 3.7, 12, 1.0, size=17, color=RGBColor(0x55,0x60,0x80), align=PP_ALIGN.CENTER, bold=True)
transition(s, "push")

# ══════════════════════════════════════════
# SLIDE 10 — Section Technologies
# ══════════════════════════════════════════
section_slide(prs, "04", "Technologies de Developpement", ACCT, PRPL)

# ══════════════════════════════════════════
# SLIDE 11 — Technologies (2 columns)
# ══════════════════════════════════════════
s = header_slide(prs, "Technologies de Developpement")

# Front-end
rect(s, 0.35, 1.6, 6.1, 5.55, DARK)
rect(s, 0.35, 1.6, 6.1, 0.15, CYAN)
txt(s, "FRONT-END", 0.6, 1.85, 5.5, 0.65, size=20, bold=True, color=CYAN)
rect(s, 0.6, 2.55, 5.6, 0.05, RGBColor(0x30,0x50,0x80))
fe = [("HTML","Structurer le contenu des pages web."),
      ("CSS","Mise en forme et esthetique des interfaces."),
      ("JavaScript","Dynamisme et interactivite navigateur."),
      ("Bootstrap","Framework responsive pour une UI moderne.")]
y = 2.7
for tech, desc in fe:
    rect(s, 0.6, y, 5.65, 0.7, RGBColor(0x10,0x20,0x50))
    txt(s, tech, 0.75, y+0.08, 1.8, 0.55, size=13, bold=True, color=CYAN)
    txt(s, desc, 2.55, y+0.08, 3.5, 0.55, size=12, color=WHITE, wrap=True)
    y += 0.82

# Back-end
rect(s, 7.0, 1.6, 6.0, 5.55, DARK)
rect(s, 7.0, 1.6, 6.0, 0.15, ACCT)
txt(s, "BACK-END", 7.2, 1.85, 5.5, 0.65, size=20, bold=True, color=ACCT)
rect(s, 7.2, 2.55, 5.5, 0.05, RGBColor(0x60,0x50,0x20))
be = [("Symfony","Framework PHP MVC robuste pour creer des applications web complexes et securisees a grande echelle."),
      ("MySQL","SGBD relationnel pour stocker et gerer toutes les donnees du cabinet (clients, dossiers, factures...).")]
y = 2.7
for tech, desc in be:
    rect(s, 7.2, y, 5.55, 1.8, RGBColor(0x10,0x20,0x50))
    txt(s, tech, 7.35, y+0.1, 2.0, 0.55, size=15, bold=True, color=ACCT)
    txt(s, desc, 7.35, y+0.65, 5.15, 1.0, size=12, color=WHITE, wrap=True)
    y += 2.1

rect(s, 6.6, 1.6, 0.08, 5.55, ACCT)
transition(s, "push")

# ══════════════════════════════════════════
# SLIDE 12 — Section Realisation
# ══════════════════════════════════════════
section_slide(prs, "05", "Realisation", CYAN, ACCT)

# ══════════════════════════════════════════
# SLIDES 13-17 — Interface screens
# ══════════════════════════════════════════
screens = [
    ("Page de Connexion (Login)",       CYAN),
    ("Tableau de Bord (Dashboard)",     ACCT),
    ("Gestion des Clients",             PRPL),
    ("Calendrier Interactif (Agenda)",  RGBColor(0x00,0xC8,0x80)),
    ("Generation de Facture PDF",       RGBColor(0xFF,0x6B,0x35)),
]
for label, color in screens:
    s = header_slide(prs, f"Realisation  |  {label}")
    rect(s, 0.5, 1.6, 12.3, 5.55, RGBColor(0xD8,0xDF,0xF5))
    rect(s, 0.5, 1.6, 12.3, 0.1,  color)
    rect(s, 0.5, 7.05, 12.3, 0.1, color)
    txt(s, f"[ Inserer ici : {label} ]",
        0.6, 3.8, 12, 0.9, size=16, color=RGBColor(0x55,0x60,0x80), align=PP_ALIGN.CENTER, bold=True)
    transition(s, "push")

# ══════════════════════════════════════════
# SLIDE 18 — Section Conclusion
# ══════════════════════════════════════════
section_slide(prs, "06", "Conclusion Generale", ACCT, PRPL)

# ══════════════════════════════════════════
# SLIDE 19 — Conclusion text
# ══════════════════════════════════════════
s = header_slide(prs, "Conclusion Generale")
rect(s, 0.4, 1.6, 12.5, 5.55, DARK)
rect(s, 0.4, 1.6, 0.12, 5.55, ACCT)
points = [
    "La plateforme LexPro repond aux besoins complexes des cabinets d'avocats.",
    "Elle regroupe gestion des clients, dossiers, agenda et facturation dans un seul espace web securise.",
    "Grâce a Symfony et au RBAC, les donnees sensibles sont protegees selon le profil de chaque utilisateur.",
    "Ce projet nous a permis de consolider nos competences techniques et de repondre a une problematique reelle.",
]
y = 2.0
for i, pt in enumerate(points):
    rect(s, 0.65, y, 0.45, 0.45, MID)
    txt(s, str(i+1), 0.65, y+0.02, 0.45, 0.42, size=16, bold=True, color=ACCT, align=PP_ALIGN.CENTER)
    txt(s, pt, 1.25, y, 11.3, 0.8, size=14, color=WHITE, wrap=True)
    y += 1.1

transition(s, "push")

# ══════════════════════════════════════════
# SLIDE 20 — Merci
# ══════════════════════════════════════════
s = blank(prs)
rect(s, 0, 0, 13.33, 7.5, DARK)
circle(s, 8.5, -1.5, 8.0, MID)
circle(s,-2.0,  4.0, 5.0, MID)
circle(s, 9.5,  1.5, 4.0, RGBColor(0x12,0x1A,0x50))
rect(s, 0, 0, 0.4, 7.5, CYAN)
rect(s, 0.4, 0, 0.08, 7.5, RGBColor(0,100,140))
rect(s, 0, 7.1, 13.33, 0.4, ACCT)
rect(s, 0.7, 2.8, 8.5, 0.06, CYAN)
txt(s, "Merci pour", 0.7, 1.1, 10, 1.4, size=58, bold=True, color=WHITE)
txt(s, "Votre Attention !", 0.7, 2.45, 10, 1.5, size=55, bold=True, color=ACCT)
txt(s, "Med Khalil Kribi", 0.7, 4.4, 9, 0.65, size=22, bold=True, color=WHITE)
txt(s, "Encadre par : M. Walid Ben Abda  |  BTS Dev. Internet",
    0.7, 5.1, 10, 0.55, size=15, color=RGBColor(0xAA,0xBB,0xDD), italic=True)
transition(s, "fade")

# ══════════════════════════════════════════
prs.save("Presentation_LexPro_V2.pptx")
print("Presentation_LexPro_V2.pptx created successfully!")
